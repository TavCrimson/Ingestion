from __future__ import annotations

from pathlib import Path

from ingestion.extractors.base import ExtractedText, ExtractorBase


class PptxExtractor(ExtractorBase):
    def extract(self, path: Path) -> ExtractedText:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(str(path))
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts = []

            # Slide title first (if present)
            if slide.shapes.title and slide.shapes.title.text.strip():
                parts.append(slide.shapes.title.text.strip())

            for shape in slide.shapes:
                # Skip the title shape — already captured above
                if shape == slide.shapes.title:
                    continue
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)

            if parts:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

        text = "\n\n".join(slides_text)

        # Extract core metadata
        metadata: dict = {}
        try:
            props = prs.core_properties
            metadata = {
                "title":    props.title or None,
                "author":   props.author or None,
                "created":  str(props.created) if props.created else None,
                "modified": str(props.modified) if props.modified else None,
                "slides":   len(prs.slides),
            }
        except Exception:
            pass

        return ExtractedText(
            text=text,
            metadata=metadata,
            pages=len(prs.slides),
        )
